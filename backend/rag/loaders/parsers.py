import os
import re
from typing import List, Dict, Any
import pypdf
import docx
from pptx import Presentation
import pandas as pd
import requests
from bs4 import BeautifulSoup

def clean_text(text: str) -> str:
    # Remove redundant whitespace and newlines
    text = re.sub(r'\s+', ' ', text)
    # Strip leading/trailing whitespaces
    return text.strip()

def parse_pdf(file_path: str) -> List[Dict[str, Any]]:
    pages = []
    with open(file_path, "rb") as f:
        reader = pypdf.PdfReader(f)
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = clean_text(text)
            if text:
                pages.append({"page": i + 1, "text": text})
    return pages

def parse_docx(file_path: str) -> List[Dict[str, Any]]:
    doc = docx.Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    pages = []
    # Group every 5 paragraphs to simulate a "page" block
    chunk_size = 5
    for i in range(0, len(paragraphs), chunk_size):
        text = clean_text(" ".join(paragraphs[i:i+chunk_size]))
        if text:
            pages.append({"page": (i // chunk_size) + 1, "text": text})
    return pages

def parse_pptx(file_path: str) -> List[Dict[str, Any]]:
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text = clean_text(" ".join(slide_text))
        if text:
            pages.append({"page": i + 1, "text": text})
    return pages

def parse_txt(file_path: str) -> List[Dict[str, Any]]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    text = clean_text(text)
    pages = []
    # Split text into segments of ~2000 characters to act as pages
    chunk_len = 2000
    for i in range(0, len(text), chunk_len):
        block = text[i:i+chunk_len]
        pages.append({"page": (i // chunk_len) + 1, "text": block})
    return pages

def parse_csv(file_path: str) -> List[Dict[str, Any]]:
    df = pd.read_csv(file_path)
    rows = []
    for idx, row in df.iterrows():
        row_str = ", ".join([f"{col}={val}" for col, val in row.items() if pd.notna(val)])
        rows.append(row_str)
    
    pages = []
    # Group every 15 rows together into chunks
    group_size = 15
    for i in range(0, len(rows), group_size):
        text = clean_text("\n".join(rows[i:i+group_size]))
        pages.append({"page": (i // group_size) + 1, "text": text})
    return pages

def parse_url(url: str) -> List[Dict[str, Any]]:
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
    response = requests.get(url, headers=headers, timeout=10)
    response.raise_for_status()
    
    soup = BeautifulSoup(response.text, "html.parser")
    # Clean the DOM of scripting/navigation clutter
    for block in soup(["script", "style", "nav", "footer", "header", "aside"]):
        block.decompose()
        
    text = soup.get_text()
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = clean_text(" ".join(chunk for chunk in chunks if chunk))
    
    pages = []
    chunk_len = 2000
    for i in range(0, len(text), chunk_len):
        block = text[i:i+chunk_len]
        pages.append({"page": (i // chunk_len) + 1, "text": block})
    return pages

def parse_document(file_path: str, file_type: str) -> List[Dict[str, Any]]:
    file_type = file_type.lower().strip(".")
    if file_type == "pdf":
        return parse_pdf(file_path)
    elif file_type in ["docx", "doc"]:
        return parse_docx(file_path)
    elif file_type in ["pptx", "ppt"]:
        return parse_pptx(file_path)
    elif file_type == "csv":
        return parse_csv(file_path)
    elif file_type == "txt":
        return parse_txt(file_path)
    elif file_type == "url":
        return parse_url(file_path)
    else:
        # Fallback to general txt parsing
        return parse_txt(file_path)
